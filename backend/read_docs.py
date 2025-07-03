"""
    read_docs.py - Helper script to read docs to GCS and prepare them for Pinecone

    Sections:
    1. Imports
    2. Global Overhead
        a. API Key and Environment Variables
        b. Database Connection
        c. Global Helper Functions
    3. File Transformation Functions
        a. extract_text_from_pdf
        b. extract_text_from_pptx
        c. chunk_text
        d. to_pinecone
    4. Main Function
"""

# --------------------------------------------------- #
# --------------------- Imports --------------------- #
# --------------------------------------------------- #

import os
import sys
import fitz  # PyMuPDF
import re
from pptx import Presentation
import io
from google.cloud import storage  # Google Cloud Storage library
from dotenv import load_dotenv
import mysql.connector

# Pinecone/LangChain imports
from pinecone import Pinecone
from pinecone import ServerlessSpec
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_pinecone import PineconeVectorStore

# ----------------------------------------------------------- #
# --------------------- Global Overhead --------------------- #
# ----------------------------------------------------------- #

os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))

# Load environment variables
load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), '.env'))

DB_HOST = os.environ.get("DB_HOST")
DB_NAME = os.environ.get("DB_NAME")
DB_USER = os.environ.get("DB_USER")
DB_PASS = os.environ.get("DB_PASS")
DB_PORT = os.environ.get("DB_PORT")

def get_db_connection():
    conn = mysql.connector.connect(
        host=DB_HOST,
        database=DB_NAME,
        user=DB_USER,
        password=DB_PASS,
        port=int(DB_PORT)
    )
    return conn

# Google Cloud Storage setup
GCS_BUCKET_NAME = 'ai-tutor-docs-bucket' 

# Initialize Google Cloud Storage client
storage_client = storage.Client()
bucket = storage_client.bucket(GCS_BUCKET_NAME)

def get_course_name(courseId):
    """
    Retrieves the course name from the database based on courseId.
    
    Args:
        courseId (int): The ID of the course.

    Returns:
        str: The name of the course.
    """
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT name FROM courses WHERE id = %s", (courseId,))
    result = cursor.fetchone()
    conn.close()
    
    if not result:
        raise ValueError(f"Course with ID {courseId} not found.")
    
    return result['name']

def get_filepath_from_db(courseId):
    """
    Retrieves the file path from the database based on courseId.
    
    Args:
        courseId (int): The ID of the course.

    Returns:
        str: The file path of the course.
    """
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT filepath FROM courses WHERE id = %s", (courseId,))
    result = cursor.fetchone()
    conn.close()
    
    if not result:
        raise ValueError(f"Course with ID {courseId} not found.")
    
    return result['filepath']

# ------------------------------------------------------------------------- #
# --------------------- File Transformation Functions --------------------- #
# ------------------------------------------------------------------------- #

# Function to read text from PDF using PyMuPDF (typed text)
def extract_text_from_pdf(pdf_bytes):
    """
    Extracts text from a PDF file.

    Args:
        pdf_bytes (bytes): The PDF file content in bytes.
    Returns:
        str: Extracted text from the PDF.
    """
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        page_text = page.get_text("text")
        # Normalize whitespace
        cleaned = re.sub(r'\s+', ' ', page_text).strip()
        text += cleaned + " /-\ "
    return text

# Function to read text from PowerPoint using python-pptx
def extract_text_from_pptx(pptx_bytes):
    """
    Extracts text from a PowerPoint presentation.

    Args:
        pptx_bytes (bytes): The PowerPoint file content in bytes.
    Returns:
        str: Extracted text from the PowerPoint.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ""
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                line = shape.text.strip()
                if line:
                    slide_text.append(line)
        if slide_text:
            slide_cleaned = " ".join(slide_text)
            slide_cleaned = re.sub(r'\s+', ' ', slide_cleaned)
            text += slide_cleaned + " /-\ "
    return text

# Function to chunk text
def chunk_text(text, chunk_size=500, chunk_overlap=100):
    """
    Splits text into smaller chunks for embedding.

    Args:
        text (str): The text to be split.
        chunk_size (int): The size of each chunk.

    Returns:
        list: List of text chunks.
    """

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size = chunk_size,
        chunk_overlap = chunk_overlap,
        separators=["\n\n", "\n", ".", "!", "?", " "]
    )

    chunks = text_splitter.split_text(text)
    return chunks

# Function to embed and store text in Pinecone
def to_pinecone(text_dict, courseId):
    """
    Embeds and stores text in Pinecone vector database.

    Args:
        text_dict (dict): Dictionary of filenames and their extracted text.
        course_name (str): The course name.
    """
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    index_name = "ai-tutor-index"

    if index_name not in pc.list_indexes().names():
        pc.create_index(
            name=index_name,
            dimension=1536,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )

    index = pc.Index(index_name)

    # Get course name from database
    course_name = get_course_name(courseId)
    # Create namespace name
    namespace = f"{course_name}_{courseId}"

    for filename, text in text_dict.items():
        # Creating chunks
        chunks = chunk_text(text)
        print(f"Embedding {len(chunks)} chunks from file: {filename}")

        vectors = embeddings.embed_documents(chunks)
        metadatas = [
            {"course_name": course_name, "filename": filename, "chunk_text": chunk}
            for chunk in chunks
        ]

        for idx, (vector, metadata) in enumerate(zip(vectors, metadatas)):
            chunk_id = f"{course_name}-{metadata['filename']}-{idx}"
            index.upsert(
                vectors=[{
                    "id": chunk_id,
                    "values": vector,
                    "metadata": metadata
                }],
                namespace=namespace
            )

    print(f"All chunks upserted to Pinecone for file {metadata['filename']}.")

# --------------------------------------------------------- #
# --------------------- Main Function --------------------- #
# --------------------------------------------------------- #

def main():
    """
    Main function to process files and store them in Pinecone.
    This function expects command-line arguments for username, courseId, proctor_id, and an optional specific file.
    """

    if len(sys.argv) < 4:
        raise ValueError("Username, Course Name, and Proctor ID are required as command-line arguments.")
    
    username = sys.argv[1]
    courseId = sys.argv[2]
    proctor_id = int(sys.argv[3])
    specific_file = sys.argv[4] if len(sys.argv) > 4 else None

    print(f"Training context for user: {username}, courseId: {courseId}, proctor ID: {proctor_id}")

    if specific_file:
        print(f"📂 Processing only file: {specific_file}")
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)

        # This is not a valid way to find filepath, should be based on course id - name is not unique
        cursor.execute("SELECT filepath FROM courses WHERE id = %s", (courseId,))
        result = cursor.fetchone()
        conn.close()
        if not courseId:
            print("Error: Course ID not found.")
            sys.exit(1)

        filepath = result['filepath']
        blob_path = filepath + specific_file
        blob = bucket.blob(blob_path)
        file_bytes = blob.download_as_bytes()

        if specific_file.endswith(".pdf"):
            text = extract_text_from_pdf(file_bytes)
        elif specific_file.endswith(".pptx"):
            text = extract_text_from_pptx(file_bytes)
        else:
            print(f"Unsupported file type: {specific_file}")
            print("Error: Unsupported file type")
            sys.exit(1)

        to_pinecone({specific_file: text}, courseId)

    else:
        print("ERROR: No specific file provided.")
        sys.exit(1)


if __name__ == "__main__":
    main()
